import pandas as pd
import sqlite3
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches


def load_data_from_excel(file_path):
    """Charge les feuilles Excel necessaires en DataFrame"""

    sheets = {
        "Vente_France": pd.read_excel(file_path, sheet_name="Vente_France"),
        "Vente_Allemagne": pd.read_excel(file_path, sheet_name="Vente_Allemagne"),
        "Vente_Pologne": pd.read_excel(file_path, sheet_name="Vente_Pologne"),
        "Cout": pd.read_excel(file_path, sheet_name="Cout"),
        "Referentiel_Produit": pd.read_excel(file_path, sheet_name="Referentiel_Produit"),
    }

    return sheets


def load_data_to_sqlite(sheets):
    """Charge les DataFrames dans une base de donnees SQLite"""

    conn = sqlite3.connect(":memory:")

    # On remplace les tables si elles existent deja
    for sheet_name, df in sheets.items():
        df.to_sql(sheet_name, conn, if_exists="replace", index=False)

    return conn


def execute_sql_query(conn, query):
    """Execute une requete SQL et retourne le resultat sous forme de DataFrame"""

    return pd.read_sql_query(query, conn)


def generate_graph_bar_from_dataframe(df, title, x_label, y_label, output_file):
    """Genere un graphique à barres à partir d'un DataFrame"""

    # print("debug head", df.head())
    print("value x_label:", x_label)
    print("value y_label:", y_label)

    plt.figure(figsize=(8, 6))

    # print("df[x_label]:", df[x_label])
    # print("df[y_label]:", df[y_label])
    print("df head IN GENERATE GRAPH:", df.head())
    print("df SHAPE", df.shape)
    print("df_x_label SHAPE", df[x_label].shape)
    print("df_y_label SHAPE", df[y_label].shape)
    print("df[x_label]:", df[x_label].head())
    print("df[y_label]:", df[y_label].head())
    print(f"df[x_label] length: {len(df[x_label])}")
    print(f"df[y_label] length: {len(df[y_label])}")


    plt.bar(df[x_label], df[y_label], color=["skyblue", "green", "red"])

    plt.title(title)
    plt.xlabel(x_label)
    plt.ylabel(y_label)
    plt.savefig(output_file)
    plt.close()


def generate_pie_chart_from_dataframe(df, labels_column, values_column, title, output_file):
    """Genere un graphique en camembert à partir d'un DataFrame"""

    plt.figure(figsize=(8, 6))
    plt.pie(df[values_column], labels=df[labels_column], autopct="%2f%%")
    plt.title(title)
    plt.savefig(output_file)
    plt.close()


def update_powerpoint(pptx_file, graphs, output_file):
    """Met à jour un fichier PowerPoint avec des graphiques"""

    ppt = Presentation(pptx_file)

    # Slide 3 - Gauche
    slide_3 = ppt.slides[2]
    slide_3.shapes.add_picture(graphs["slide3_left"]["graph"], Inches(1), Inches(2), width=Inches(5), height=Inches(5))
    textbox = slide_3.shapes.add_textbox(Inches(1), Inches(6), width=Inches(8), height=Inches(1))
    textbox.text = graphs["slide3_left"]["text"]

    # Slide 3 - Droite
    slide_3.shapes.add_picture(graphs["slide3_right"]["graph"], Inches(6), Inches(2), width=Inches(5), height=Inches(5))
    textbox = slide_3.shapes.add_textbox(Inches(1), Inches(6), width=Inches(8), height=Inches(1))
    textbox.text = graphs["slide3_right"]["text"]

    # # Slide 4 - Gauche
    # slide_4 = ppt.slides[3]
    # slide_4.shapes.add_picture(graphs["slide4_left"]["graph"], Inches(1), Inches(2), width=Inches(5), height=Inches(5))
    # textbox = slide_4.shapes.add_textbox(Inches(1), Inches(6), width=Inches(8), height=Inches(1))
    # textbox.text = graphs["slide4_left"]["text"]

    # # Slide 4 - Droite
    # slide_4.shapes.add_picture(graphs["slide4_right"]["graph"], Inches(6), Inches(2), width=Inches(5), height=Inches(5))    
    # textbox = slide_4.shapes.add_textbox(Inches(1), Inches(6), width=Inches(8), height=Inches(1))
    # textbox.text = graphs["slide4_right"]["text"]

    ppt.save(output_file)


def main():

    # Chargement des données
    excel_file = "QuadraticData_Test_Données.xlsx"
    ppt_file = "QuadraticData - Test recrutement.pptx"
    ppt_output = "QuadraticData - Test recrutement - Resultat.pptx"

    # Chargement des feuilles Excel et Sqlite
    sheets = load_data_from_excel(excel_file)
    conn = load_data_to_sqlite(sheets)

    graphs = {}

    # Slide 3 - Gauche : CA par pays
    query_ca_by_country = """
    SELECT Pays, SUM(CA) AS Total_CA
    FROM (
        SELECT 'France' AS Pays, Ca FROM Vente_France
        UNION ALL
        SELECT 'Allemagne' AS Pays, Ca FROM Vente_Allemagne
        UNION ALL
        SELECT 'Pologne' AS Pays, Ca FROM Vente_Pologne
    )
    GROUP BY Pays
    ORDER BY Total_CA DESC
    """

    df_ca = execute_sql_query(conn, query_ca_by_country)
    generate_graph_bar_from_dataframe(df_ca, "Chiffre d'affaires par Pays", "Pays", "Total_CA", "ca_by_country.png")
    top_country = df_ca.iloc[0]
    total_ca = df_ca["Total_CA"].sum()
    ca_text = f"Pays avec le CA le plus élevé : {top_country['Pays']} avec {top_country['Total_CA']} ({(top_country['Total_CA'] / total_ca) * 100:.2f}%)"
    graphs["slide3_left"] = {"graph": "ca_by_country.png", "text": ca_text}

    # Slide 3 - Droite : Évolution du CA
    query_ca_evolution = """
    SELECT Pays,
              SUM(CASE WHEN Annee = 2020 THEN CA ELSE 0 END) AS CA_2020,
              SUM(CASE WHEN Annee = 2019 THEN CA ELSE 0 END) AS CA_2019,
              SUM(CASE WHEN Annee = 2020 THEN CA ELSE 0 END) - SUM(CASE WHEN Annee = 2019 THEN CA ELSE 0 END) AS Rate
    FROM (
        SELECT 'France' AS Pays, Annee, Ca FROM Vente_France
        UNION ALL
        SELECT 'Allemagne' AS Pays, Annee, Ca FROM Vente_Allemagne
        UNION ALL
        SELECT 'Pologne' AS Pays, Annee, Ca FROM Vente_Pologne
    )
    GROUP BY Pays
    ORDER BY Rate DESC
    """
    df_evolution = execute_sql_query(conn, query_ca_evolution)

    generate_graph_bar_from_dataframe(df_evolution, "Pays", "Rate", "Rate", "ca_evolution.png")
    top_evolution = df_evolution.iloc[0]
    evolution_text = f"Pays avec la plus forte évolution du CA : {top_evolution['Pays']} avec un taux de {top_evolution['Rate']}"
    graphs["slide3_right"] = {"graph": "ca_evolution.png", "text": evolution_text}

    # # Slide 4 - Gauche : Coût par produit
    # query_cost_per_product = """
    # SELECT v.Produit,
    #             SUM(v.CA - c.Cout_France) AS Marge_France,
    #             SUM(v.CA - c.Cout_Allemagne) AS Marge_Allemagne,
    #             SUM(v.CA - c.Cout_Pologne) AS Marge_Pologne
    # FROM (
    #     SELECT Produit, CA FROM Vente_France
    #     UNION ALL
    #     SELECT Produit, CA FROM Vente_Allemagne
    #     UNION ALL
    #     SELECT Produit, CA FROM Vente_Pologne
    # ) v
    # JOIN Cout c ON v.Produit = c.Produit
    # GROUP BY v.Produit
    # ORDER BY Marge_France DESC
    # """
    # df_marge = execute_sql_query(conn, query_cost_per_product)
    # df_marge = df_marge.rename(columns={"Marge_France": "Marge",
    #                                     "Marge_Allemagne": "Marge",
    #                                     "Marge_Pologne": "Marge"})
    
    # generate_graph_bar_from_dataframe(df_marge, "Marge par produit", "Produit", "Marge", "marge_per_product.png")
    # top_product = df_marge.iloc[0]
    # total_marge = df_marge["Marge"].sum()
    # marge_text = f"Produit avec la plus grande marge : {top_product['Produit']} avec {top_product['Marge']} ({(top_product['Marge'] / total_marge) * 100:.2f}%)"
    # graphs["slide4_left"] = {"graph": "marge_per_product.png", "text": marge_text}

    # # Slide 4 - Droite : Répartition des coûts
    # query_cost_distribution = """
    # SELECT Pays,
    #             SUM(CA - Cout_France) AS Marge_France,
    #             SUM(CA - Cout_Allemagne) AS Marge_Allemagne,
    #             SUM(CA - Cout_Pologne) AS Marge_Pologne
    # FROM (
    #     SELECT 'France' AS Pays, Produit, CA FROM Vente_France
    #     UNION ALL
    #     SELECT 'Allemagne' AS Pays, Produit, CA FROM Vente_Allemagne
    #     UNION ALL
    #     SELECT 'Pologne' AS Pays, Produit, CA FROM Vente_Pologne
    # ) v
    # JOIN Cout c ON v.Produit = c.Produit
    # WHERE v.Produit = '{top_product['Produit']}'
    # GROUP BY Pays
    # ORDER BY Marge_France DESC
    # """
    # df_distribution = execute_sql_query(conn, query_cost_distribution)
    # generate_pie_chart_from_dataframe(df_distribution, "Pays", "Marge", "Répartition des coûts", "cost_distribution.png")
    # top_country_cost = df_distribution.iloc[0]
    # cost_text = f"Pays avec la plus grande marge : {top_country_cost['Pays']} avec {top_country_cost['Marge']} ({(top_country_cost['Marge'] / total_marge) * 100:.2f}%)"
    # graphs["slide4_right"] = {"graph": "cost_distribution.png", "text": cost_text}

    update_powerpoint(ppt_file, graphs, ppt_output)

    print(f"Le fichier PowerPoint a été généré avec succès : {ppt_output}")


if __name__ == "__main__":
    main()
